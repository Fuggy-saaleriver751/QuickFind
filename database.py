"""
QuickFind — Database, Indexer & File Watcher
SQLite FTS5 · 3-Phase Indexing · Parallel Scanning · Real-time Watching

Phase 1: Fast metadata scan (name, size, date) — search ready in seconds
Phase 2: Content indexing (text/rich documents) — FTS enrichment in background
Phase 3: Hash computation (SHA-256) — duplicate detection in background

Optimizations:
- Directory deduplication (dirs table) — paths stored once, referenced by id
- FTS indexes name + extension + folder_name + content only
- ThreadPoolExecutor for parallel directory scanning
- Parallel drive detection (A:-Z:, dynamic)
- WAL mode + batch inserts
"""

import sqlite3
import os
import re
import time
import threading
import json
import hashlib
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler


TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".log", ".cfg", ".ini", ".conf",
    ".py", ".pyw", ".js", ".ts", ".jsx", ".tsx", ".vue", ".svelte",
    ".html", ".htm", ".css", ".scss", ".less",
    ".java", ".kt", ".scala",
    ".c", ".cpp", ".cc", ".h", ".hpp",
    ".cs", ".go", ".rs", ".rb", ".php",
    ".sh", ".bash", ".bat", ".cmd", ".ps1",
    ".json", ".xml", ".yaml", ".yml", ".toml",
    ".csv", ".sql", ".env", ".gitignore",
    ".tex", ".bib", ".srt",
}

RICH_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx", ".rtf", ".epub",
}

ALL_CONTENT_EXTENSIONS = TEXT_EXTENSIONS | RICH_EXTENSIONS

MAX_RICH_FILE_SIZE = 50_000_000

# ── Content indexing presets ──
PRESETS = {
    "minimal":  {"text_bytes": 128,   "rich_chars": 128,   "label": "Minimal (~50-80 MB) *"},
    "standard": {"text_bytes": 512,   "rich_chars": 512,   "label": "Standard (~80-150 MB) *"},
    "deep":     {"text_bytes": 4096,  "rich_chars": 4096,  "label": "Deep (~200-500 MB) *"},
    "maximum":  {"text_bytes": 0,     "rich_chars": 0,     "label": "Maximum (~500 MB+) *"},
}

import sys as _sys
if getattr(_sys, 'frozen', False):
    _SCRIPT_DIR = os.path.dirname(_sys.executable)
else:
    _SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DB_DIR = os.path.join(_SCRIPT_DIR, "QuickFind_Index")
DB_PATH = os.path.join(DB_DIR, "index.db")
CONFIG_PATH = os.path.join(DB_DIR, "config.json")


def _load_config():
    try:
        with open(CONFIG_PATH, "r") as f:
            return json.load(f)
    except Exception:
        return {"preset": "minimal"}

def _save_config(cfg):
    os.makedirs(os.path.dirname(CONFIG_PATH), exist_ok=True)
    with open(CONFIG_PATH, "w") as f:
        json.dump(cfg, f, indent=2)

def get_preset_name():
    return _load_config().get("preset", "minimal")

def set_preset_name(name):
    cfg = _load_config()
    cfg["preset"] = name
    _save_config(cfg)
    _reset_limits_cache()

_limits_cache = None

def _get_limits():
    global _limits_cache
    if _limits_cache is None:
        preset = get_preset_name()
        p = PRESETS.get(preset, PRESETS["minimal"])
        _limits_cache = (p["text_bytes"], p["rich_chars"])
    return _limits_cache

def _reset_limits_cache():
    global _limits_cache
    _limits_cache = None


# ══════════════════════════════════════════════════════════════
#  RICH DOCUMENT READERS
# ══════════════════════════════════════════════════════════════

_READ_LIMIT = 8192

def _read_pdf(path):
    try:
        import fitz
        doc = fitz.open(path)
        parts = []
        for page in doc:
            parts.append(page.get_text())
            if sum(len(t) for t in parts) > _READ_LIMIT:
                break
        doc.close()
        return "".join(parts)[:_READ_LIMIT]
    except Exception:
        return ""

def _read_docx(path):
    try:
        from docx import Document
        doc = Document(path)
        parts, total = [], 0
        for para in doc.paragraphs:
            parts.append(para.text)
            total += len(para.text)
            if total > _READ_LIMIT: break
        return "\n".join(parts)[:_READ_LIMIT]
    except Exception:
        return ""

def _read_xlsx(path):
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts, total = [], 0
        for sheet in wb.sheetnames:
            for row in wb[sheet].iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    line = " ".join(vals)
                    parts.append(line)
                    total += len(line)
                    if total > _READ_LIMIT: break
            if total > _READ_LIMIT: break
        wb.close()
        return "\n".join(parts)[:_READ_LIMIT]
    except Exception:
        return ""

def _read_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts, total = [], 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        parts.append(para.text)
                        total += len(para.text)
                        if total > _READ_LIMIT: break
                if total > _READ_LIMIT: break
            if total > _READ_LIMIT: break
        return "\n".join(parts)[:_READ_LIMIT]
    except Exception:
        return ""

def _read_rtf(path):
    try:
        from striprtf.striprtf import rtf_to_text
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read(_READ_LIMIT * 4)
        return rtf_to_text(raw)[:_READ_LIMIT]
    except Exception:
        return ""

def _read_epub(path):
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(path, options={"ignore_ncx": True})
        parts, total = [], 0
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator=" ", strip=True)
            parts.append(text)
            total += len(text)
            if total > _READ_LIMIT: break
        return "\n".join(parts)[:_READ_LIMIT]
    except Exception:
        return ""

RICH_READERS = {
    ".pdf": _read_pdf, ".docx": _read_docx, ".xlsx": _read_xlsx,
    ".pptx": _read_pptx, ".rtf": _read_rtf, ".epub": _read_epub,
}


# ══════════════════════════════════════════════════════════════
#  DATABASE — with directory deduplication
# ══════════════════════════════════════════════════════════════

def _regexp(pattern, value):
    """SQLite REGEXP function for regex search."""
    if value is None:
        return False
    try:
        return bool(re.search(pattern, value, re.IGNORECASE))
    except re.error:
        return False


class FileDatabase:
    def __init__(self):
        os.makedirs(DB_DIR, exist_ok=True)
        self.conn = sqlite3.connect(DB_PATH, check_same_thread=False)
        self.conn.execute("PRAGMA journal_mode=WAL")
        self.conn.execute("PRAGMA synchronous=NORMAL")
        self.conn.execute("PRAGMA cache_size=-64000")
        self.conn.execute("PRAGMA wal_autocheckpoint=200")
        self.conn.execute("PRAGMA page_size=4096")
        self.conn.execute("PRAGMA mmap_size=268435456")  # 256MB mmap for faster reads
        self.conn.create_function("REGEXP", 2, _regexp)
        self.lock = threading.Lock()
        self._dir_cache = {}
        self._create_tables()

    def _create_tables(self):
        with self.lock:
            self.conn.executescript("""
                CREATE TABLE IF NOT EXISTS dirs (
                    id INTEGER PRIMARY KEY,
                    path TEXT NOT NULL UNIQUE
                );

                CREATE TABLE IF NOT EXISTS files (
                    id INTEGER PRIMARY KEY,
                    name TEXT NOT NULL,
                    dir_id INTEGER NOT NULL REFERENCES dirs(id),
                    extension TEXT,
                    size INTEGER DEFAULT 0,
                    modified REAL DEFAULT 0,
                    folder_name TEXT DEFAULT '',
                    content_text TEXT DEFAULT '',
                    UNIQUE(dir_id, name)
                );

                CREATE VIRTUAL TABLE IF NOT EXISTS files_fts USING fts5(
                    name,
                    extension,
                    folder_name,
                    content_text,
                    content='files',
                    content_rowid='id',
                    tokenize='unicode61 remove_diacritics 2'
                );

                CREATE TRIGGER IF NOT EXISTS files_ai AFTER INSERT ON files BEGIN
                    INSERT INTO files_fts(rowid, name, extension, folder_name, content_text)
                    VALUES (new.id, new.name, new.extension, new.folder_name, new.content_text);
                END;

                CREATE TRIGGER IF NOT EXISTS files_ad AFTER DELETE ON files BEGIN
                    INSERT INTO files_fts(files_fts, rowid, name, extension, folder_name, content_text)
                    VALUES ('delete', old.id, old.name, old.extension, old.folder_name, old.content_text);
                END;

                CREATE TRIGGER IF NOT EXISTS files_au AFTER UPDATE ON files BEGIN
                    INSERT INTO files_fts(files_fts, rowid, name, extension, folder_name, content_text)
                    VALUES ('delete', old.id, old.name, old.extension, old.folder_name, old.content_text);
                    INSERT INTO files_fts(rowid, name, extension, folder_name, content_text)
                    VALUES (new.id, new.name, new.extension, new.folder_name, new.content_text);
                END;

                CREATE TABLE IF NOT EXISTS meta (
                    key TEXT PRIMARY KEY,
                    value TEXT
                );

                CREATE INDEX IF NOT EXISTS idx_files_ext ON files(extension);
                CREATE INDEX IF NOT EXISTS idx_files_dir ON files(dir_id);
                CREATE INDEX IF NOT EXISTS idx_files_modified ON files(modified);
                CREATE INDEX IF NOT EXISTS idx_files_size ON files(size);
            """)
            self.conn.commit()

            # Migration: add content_hash column if missing
            try:
                self.conn.execute("SELECT content_hash FROM files LIMIT 0")
            except sqlite3.OperationalError:
                self.conn.execute("ALTER TABLE files ADD COLUMN content_hash TEXT DEFAULT ''")
                self.conn.commit()
            self.conn.execute("CREATE INDEX IF NOT EXISTS idx_files_hash ON files(content_hash)")
            self.conn.commit()

    def _get_dir_id(self, dir_path):
        if dir_path in self._dir_cache:
            return self._dir_cache[dir_path]
        cur = self.conn.execute("SELECT id FROM dirs WHERE path=?", (dir_path,))
        row = cur.fetchone()
        if row:
            self._dir_cache[dir_path] = row[0]
            return row[0]
        self.conn.execute("INSERT OR IGNORE INTO dirs(path) VALUES(?)", (dir_path,))
        self.conn.commit()
        cur = self.conn.execute("SELECT id FROM dirs WHERE path=?", (dir_path,))
        row = cur.fetchone()
        self._dir_cache[dir_path] = row[0]
        return row[0]

    def get_meta(self, key):
        cur = self.conn.execute("SELECT value FROM meta WHERE key=?", (key,))
        row = cur.fetchone()
        return row[0] if row else None

    def set_meta(self, key, value):
        with self.lock:
            self.conn.execute("INSERT OR REPLACE INTO meta(key,value) VALUES(?,?)", (key, str(value)))
            self.conn.commit()

    def get_file_count(self):
        cur = self.conn.execute("SELECT COUNT(*) FROM files")
        return cur.fetchone()[0]

    def clear(self):
        with self.lock:
            self.conn.executescript("""
                DELETE FROM files;
                DELETE FROM dirs;
                INSERT INTO files_fts(files_fts) VALUES('rebuild');
            """)
            self.conn.commit()
            self._dir_cache.clear()

    def insert_batch(self, file_list):
        """file_list: [(name, dir_path, ext, size, modified, folder_name, content_text, content_hash), ...]
        content_text and content_hash are optional (can be 7 or 8 element tuples)."""
        with self.lock:
            rows = []
            for item in file_list:
                name, dir_path, ext, size, modified, folder_name = item[:6]
                content_text = item[6] if len(item) > 6 else ""
                content_hash = item[7] if len(item) > 7 else ""
                dir_id = self._get_dir_id(dir_path)
                rows.append((name, dir_id, ext, size, modified, folder_name, content_text, content_hash))
            self.conn.executemany(
                "INSERT OR IGNORE INTO files(name, dir_id, extension, size, modified, folder_name, content_text, content_hash) "
                "VALUES(?, ?, ?, ?, ?, ?, ?, ?)",
                rows
            )
            self.conn.commit()
            self.conn.execute("PRAGMA wal_checkpoint(PASSIVE)")

    def upsert_batch(self, file_list):
        with self.lock:
            rows = []
            for item in file_list:
                name, dir_path, ext, size, modified, folder_name = item[:6]
                content_text = item[6] if len(item) > 6 else ""
                content_hash = item[7] if len(item) > 7 else ""
                dir_id = self._get_dir_id(dir_path)
                rows.append((name, dir_id, ext, size, modified, folder_name, content_text, content_hash))
            self.conn.executemany(
                "INSERT INTO files(name, dir_id, extension, size, modified, folder_name, content_text, content_hash) "
                "VALUES(?, ?, ?, ?, ?, ?, ?, ?) "
                "ON CONFLICT(dir_id, name) DO UPDATE SET "
                "extension=excluded.extension, size=excluded.size, "
                "modified=excluded.modified, folder_name=excluded.folder_name, "
                "content_text=CASE WHEN excluded.content_text != '' THEN excluded.content_text ELSE files.content_text END, "
                "content_hash=CASE WHEN excluded.content_hash != '' THEN excluded.content_hash ELSE files.content_hash END",
                rows
            )
            self.conn.commit()
            self.conn.execute("PRAGMA wal_checkpoint(PASSIVE)")

    def update_content_batch(self, updates):
        """Update content_text for existing files. updates: [(content_text, dir_path, name), ...]"""
        with self.lock:
            rows = []
            for content_text, dir_path, name in updates:
                dir_id = self._dir_cache.get(dir_path)
                if dir_id is None:
                    cur = self.conn.execute("SELECT id FROM dirs WHERE path=?", (dir_path,))
                    row = cur.fetchone()
                    if not row:
                        continue
                    dir_id = row[0]
                    self._dir_cache[dir_path] = dir_id
                rows.append((content_text, dir_id, name))
            if rows:
                self.conn.executemany(
                    "UPDATE files SET content_text=? WHERE dir_id=? AND name=?", rows)
                self.conn.commit()

    def update_hash_batch(self, updates):
        """Update content_hash for existing files. updates: [(content_hash, dir_path, name), ...]"""
        with self.lock:
            rows = []
            for content_hash, dir_path, name in updates:
                dir_id = self._dir_cache.get(dir_path)
                if dir_id is None:
                    cur = self.conn.execute("SELECT id FROM dirs WHERE path=?", (dir_path,))
                    row = cur.fetchone()
                    if not row:
                        continue
                    dir_id = row[0]
                    self._dir_cache[dir_path] = dir_id
                rows.append((content_hash, dir_id, name))
            if rows:
                self.conn.executemany(
                    "UPDATE files SET content_hash=? WHERE dir_id=? AND name=?", rows)
                self.conn.commit()

    def get_files_needing_content(self, limit=500):
        """Get files that need content indexing (Phase 2).
        Only returns files with content_text exactly '' (not ' ' which means attempted)."""
        placeholders = ",".join("?" * len(ALL_CONTENT_EXTENSIONS))
        cur = self.conn.execute(f"""
            SELECT f.name, d.path, f.extension, f.size
            FROM files f JOIN dirs d ON d.id = f.dir_id
            WHERE f.content_text = '' AND f.extension IN ({placeholders})
            AND f.size > 0 AND f.size < ?
            ORDER BY f.size ASC
            LIMIT ?
        """, list(ALL_CONTENT_EXTENSIONS) + [MAX_RICH_FILE_SIZE, limit])
        return cur.fetchall()

    def get_files_needing_hash(self, limit=500):
        """Get files that need hash computation (Phase 3)."""
        cur = self.conn.execute("""
            SELECT f.name, d.path, f.size
            FROM files f JOIN dirs d ON d.id = f.dir_id
            WHERE f.content_hash = '' AND f.size > 0
            LIMIT ?
        """, (limit,))
        return cur.fetchall()

    def upsert_file(self, name, dir_path, ext, size, modified, folder_name="", content_text="", content_hash=""):
        with self.lock:
            dir_id = self._get_dir_id(dir_path)
            self.conn.execute(
                "INSERT INTO files(name, dir_id, extension, size, modified, folder_name, content_text, content_hash) "
                "VALUES(?, ?, ?, ?, ?, ?, ?, ?) "
                "ON CONFLICT(dir_id, name) DO UPDATE SET "
                "extension=excluded.extension, size=excluded.size, "
                "modified=excluded.modified, folder_name=excluded.folder_name, "
                "content_text=excluded.content_text, content_hash=excluded.content_hash",
                (name, dir_id, ext, size, modified, folder_name, content_text, content_hash)
            )
            self.conn.commit()

    def delete_by_path(self, full_path):
        dir_path = os.path.dirname(full_path)
        name = os.path.basename(full_path)
        with self.lock:
            dir_id = self._dir_cache.get(dir_path)
            if dir_id is None:
                cur = self.conn.execute("SELECT id FROM dirs WHERE path=?", (dir_path,))
                row = cur.fetchone()
                if not row:
                    return
                dir_id = row[0]
            self.conn.execute("DELETE FROM files WHERE dir_id=? AND name=?", (dir_id, name))
            self.conn.commit()

    def delete_by_dir_prefix(self, dir_prefix):
        with self.lock:
            self.conn.execute("""
                DELETE FROM files WHERE dir_id IN (
                    SELECT id FROM dirs WHERE path = ? OR path LIKE ?
                )
            """, (dir_prefix, dir_prefix + "\\%"))
            self.conn.commit()

    # ─── Search ──────────────────────────────────────────

    SORT_OPTIONS = {
        "relevance": "rank",
        "name_asc":  "f.name ASC",
        "name_desc": "f.name DESC",
        "size_asc":  "f.size ASC",
        "size_desc": "f.size DESC",
        "date_new":  "f.modified DESC",
        "date_old":  "f.modified ASC",
    }

    def search(self, query, limit=200, ext_filter=None, sort="relevance",
               min_size=None, max_size=None, date_from=None, date_to=None,
               folder_filter=None, regex_pattern=None, content_filter=None):
        """Advanced search with multiple filter types."""
        query = query.strip() if query else ""

        # Parse inline operators (ext:, hash:) for backward compat
        parsed_exts = []
        search_terms = []
        hash_query = None
        for term in (query.split() if query else []):
            low = term.lower()
            if low.startswith("ext:"):
                for e in term[4:].split(","):
                    e = e.strip().lower()
                    if not e.startswith("."):
                        e = "." + e
                    parsed_exts.append(e)
            elif low.startswith("hash:"):
                hash_query = term[5:].strip()
            else:
                search_terms.append(term)

        if hash_query:
            return self._search_by_hash(hash_query, limit)

        # Regex search mode
        if regex_pattern:
            return self._search_regex(regex_pattern, limit, ext_filter or parsed_exts,
                                       min_size, max_size, date_from, date_to, folder_filter)

        # Combine extensions
        all_exts = parsed_exts
        if ext_filter:
            all_exts = ext_filter if not parsed_exts else parsed_exts

        if not search_terms and not all_exts and not content_filter and \
           min_size is None and max_size is None and date_from is None and folder_filter is None:
            return []

        # Build extra WHERE clauses
        extra_clauses = []
        extra_params = []
        if all_exts:
            placeholders = ",".join("?" * len(all_exts))
            extra_clauses.append(f"f.extension IN ({placeholders})")
            extra_params.extend(all_exts)
        if min_size is not None:
            extra_clauses.append("f.size >= ?")
            extra_params.append(min_size)
        if max_size is not None:
            extra_clauses.append("f.size <= ?")
            extra_params.append(max_size)
        if date_from is not None:
            extra_clauses.append("f.modified >= ?")
            extra_params.append(date_from)
        if date_to is not None:
            extra_clauses.append("f.modified <= ?")
            extra_params.append(date_to)
        if folder_filter:
            extra_clauses.append("f.folder_name LIKE ?")
            extra_params.append(f"%{folder_filter}%")
        if content_filter:
            extra_clauses.append("f.content_text LIKE ?")
            extra_params.append(f"%{content_filter}%")

        extra_where = ""
        if extra_clauses:
            extra_where = "AND " + " AND ".join(extra_clauses)

        order = self.SORT_OPTIONS.get(sort, "rank")

        if search_terms:
            fts_query = " AND ".join(f'"{t}"*' for t in search_terms if t)
            try:
                cur = self.conn.execute(f"""
                    SELECT f.name, d.path || '\\' || f.name, f.extension, f.size, f.modified, 0,
                           bm25(files_fts, 10.0, 5.0, 3.0, 2.0) as rank
                    FROM files_fts
                    JOIN files f ON f.id = files_fts.rowid
                    JOIN dirs d ON d.id = f.dir_id
                    WHERE files_fts MATCH ? {extra_where}
                    ORDER BY {order}
                    LIMIT ?
                """, [fts_query] + extra_params + [limit])
                return cur.fetchall()
            except Exception:
                like = f"%{search_terms[0]}%"
                fallback_order = order if order != "rank" else "f.name ASC"
                cur = self.conn.execute(f"""
                    SELECT f.name, d.path || '\\' || f.name, f.extension, f.size, f.modified, 0, 0
                    FROM files f
                    JOIN dirs d ON d.id = f.dir_id
                    WHERE f.name LIKE ? {extra_where}
                    ORDER BY {fallback_order}
                    LIMIT ?
                """, [like] + extra_params + [limit])
                return cur.fetchall()
        else:
            no_fts_order = order if order != "rank" else "f.modified DESC"
            where = "WHERE 1=1 " + extra_where if extra_where else "WHERE 1=1"
            cur = self.conn.execute(f"""
                SELECT f.name, d.path || '\\' || f.name, f.extension, f.size, f.modified, 0, 0
                FROM files f
                JOIN dirs d ON d.id = f.dir_id
                {where}
                ORDER BY {no_fts_order}
                LIMIT ?
            """, extra_params + [limit])
            return cur.fetchall()

    def _search_regex(self, pattern, limit, ext_filter, min_size, max_size, date_from, date_to, folder_filter):
        """Search using REGEXP on file names."""
        clauses = ["f.name REGEXP ?"]
        params = [pattern]
        if ext_filter:
            placeholders = ",".join("?" * len(ext_filter))
            clauses.append(f"f.extension IN ({placeholders})")
            params.extend(ext_filter)
        if min_size is not None:
            clauses.append("f.size >= ?")
            params.append(min_size)
        if max_size is not None:
            clauses.append("f.size <= ?")
            params.append(max_size)
        if date_from is not None:
            clauses.append("f.modified >= ?")
            params.append(date_from)
        if date_to is not None:
            clauses.append("f.modified <= ?")
            params.append(date_to)
        if folder_filter:
            clauses.append("f.folder_name LIKE ?")
            params.append(f"%{folder_filter}%")
        where = " AND ".join(clauses)
        cur = self.conn.execute(f"""
            SELECT f.name, d.path || '\\' || f.name, f.extension, f.size, f.modified, 0, 0
            FROM files f
            JOIN dirs d ON d.id = f.dir_id
            WHERE {where}
            ORDER BY f.name ASC
            LIMIT ?
        """, params + [limit])
        return cur.fetchall()

    def get_file_hash(self, full_path):
        dir_path = os.path.dirname(full_path)
        name = os.path.basename(full_path)
        cur = self.conn.execute("""
            SELECT f.content_hash FROM files f
            JOIN dirs d ON d.id = f.dir_id
            WHERE d.path = ? AND f.name = ?
        """, (dir_path, name))
        row = cur.fetchone()
        return row[0] if row else ""

    def _search_by_hash(self, hash_prefix, limit=200):
        cur = self.conn.execute("""
            SELECT f.name, d.path || '\\' || f.name, f.extension, f.size, f.modified, 0, 0
            FROM files f
            JOIN dirs d ON d.id = f.dir_id
            WHERE f.content_hash LIKE ?
            ORDER BY f.name ASC
            LIMIT ?
        """, (hash_prefix + "%", limit))
        return cur.fetchall()

    def find_duplicates(self, limit=500, min_size=1):
        cur = self.conn.execute("""
            SELECT f.content_hash, f.name, d.path || '\\' || f.name, f.extension, f.size, f.modified,
                   COUNT(*) OVER (PARTITION BY f.content_hash) as dup_count
            FROM files f
            JOIN dirs d ON d.id = f.dir_id
            WHERE f.content_hash != '' AND f.size >= ?
            AND f.content_hash IN (
                SELECT content_hash FROM files
                WHERE content_hash != '' AND size >= ?
                GROUP BY content_hash HAVING COUNT(*) > 1
            )
            ORDER BY f.size DESC, f.content_hash, f.name
            LIMIT ?
        """, (min_size, min_size, limit))
        return cur.fetchall()

    def find_similar(self, name, limit=20):
        """Find files with similar names using FTS."""
        # Split name into words
        base = Path(name).stem
        words = re.split(r'[_\-.\s]+', base)
        words = [w for w in words if len(w) > 2]
        if not words:
            return []
        fts_query = " OR ".join(f'"{w}"*' for w in words[:5])
        try:
            cur = self.conn.execute("""
                SELECT f.name, d.path || '\\' || f.name, f.extension, f.size, f.modified, 0,
                       bm25(files_fts) as rank
                FROM files_fts
                JOIN files f ON f.id = files_fts.rowid
                JOIN dirs d ON d.id = f.dir_id
                WHERE files_fts MATCH ? AND f.name != ?
                ORDER BY rank
                LIMIT ?
            """, (fts_query, name, limit))
            return cur.fetchall()
        except Exception:
            return []

    def close(self):
        self.conn.close()

    def get_db_size_mb(self):
        total = 0
        for f in os.listdir(DB_DIR):
            fp = os.path.join(DB_DIR, f)
            if os.path.isfile(fp) and f.startswith("index.db"):
                total += os.path.getsize(fp)
        return total / (1024 * 1024)


# ══════════════════════════════════════════════════════════════
#  FILE INDEXER — 3-Phase, Parallel
# ══════════════════════════════════════════════════════════════

def _detect_drives():
    """Dynamically detect all available drives (A:-Z:)."""
    drives = []
    for letter in "ABCDEFGHIJKLMNOPQRSTUVWXYZ":
        drive = f"{letter}:\\"
        if os.path.exists(drive):
            drives.append(drive)
    return drives


class FileIndexer:
    SKIP_DIRS = {
        '$Recycle.Bin', '$WinREAgent', 'System Volume Information',
        'Windows', 'ProgramData', 'Recovery', 'PerfLogs',
        '.git', 'node_modules', '__pycache__', '.venv', 'venv',
        'AppData', '.cache', '.npm', '.nuget', 'packages',
        'Windows.old', 'Config.Msi', 'QuickFind_Index',
    }

    SKIP_EXTENSIONS = {
        '.tmp', '.log', '.bak', '.dmp', '.etl', '.evtx',
        '.dat', '.db-journal', '.db-shm', '.db-wal',
    }

    def __init__(self, db: FileDatabase, progress_callback=None, status_callback=None,
                 phase_callback=None):
        """
        phase_callback(phase: int, status: str) — called when phase changes:
            1 = metadata scan started/done
            2 = content indexing started/done
            3 = hash computation started/done
        """
        self.db = db
        self.progress_callback = progress_callback
        self.status_callback = status_callback
        self.phase_callback = phase_callback
        self._stop_event = threading.Event()
        self._thread = None
        self.total_indexed = 0
        self._incremental_since = 0
        self._is_reindex = True
        self.current_phase = 0

    def start(self, drives=None, reindex=False):
        if self._thread and self._thread.is_alive():
            return
        if drives is None:
            drives = _detect_drives()
        self._stop_event.clear()
        self._thread = threading.Thread(target=self._index_worker, args=(drives, reindex), daemon=True)
        self._thread.start()

    def stop(self):
        self._stop_event.set()
        if self._thread:
            self._thread.join(timeout=5)

    def is_running(self):
        return self._thread is not None and self._thread.is_alive()

    def _status(self, msg):
        if self.status_callback:
            self.status_callback(msg)

    def _progress(self, count):
        if self.progress_callback:
            self.progress_callback(count)

    def _phase(self, phase, status):
        self.current_phase = phase
        if self.phase_callback:
            self.phase_callback(phase, status)

    @staticmethod
    def compute_file_hash(path, size=0, max_bytes=8192):
        try:
            h = hashlib.sha256()
            h.update(str(size).encode())
            with open(path, 'rb') as f:
                data = f.read(max_bytes)
                if data:
                    h.update(data)
            return h.hexdigest()[:32]
        except (PermissionError, OSError):
            return ""

    @staticmethod
    def read_content(path, ext, size=0):
        text_bytes, rich_chars = _get_limits()
        if text_bytes == 0: text_bytes = 1_000_000
        if rich_chars == 0: rich_chars = 1_000_000
        if ext in RICH_EXTENSIONS:
            if size > MAX_RICH_FILE_SIZE:
                return ""
            reader = RICH_READERS.get(ext)
            if reader:
                try:
                    return reader(path)[:rich_chars]
                except Exception:
                    return ""
        elif ext in TEXT_EXTENSIONS:
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read(text_bytes)
            except (PermissionError, OSError):
                return ""
        return ""

    # ─── Main Worker ─────────────────────────────────────

    def _index_worker(self, drives, reindex):
        start_time = time.time()
        _reset_limits_cache()

        if reindex:
            self._status("Clearing old index...")
            self.db.clear()
            self._incremental_since = 0
        else:
            lt = self.db.get_meta("last_index_time")
            if lt:
                try:
                    self._incremental_since = datetime.fromisoformat(lt).timestamp() - 60
                except Exception:
                    self._incremental_since = 0
            else:
                self._incremental_since = 0

        self._is_reindex = reindex
        self.total_indexed = 0

        # ═══ PHASE 1: Fast metadata scan (parallel per drive) ═══
        self._phase(1, "starting")
        self._status("Phase 1: Scanning file metadata...")
        self._phase1_scan(drives)
        p1_time = time.time() - start_time

        if self._stop_event.is_set():
            return

        count = self.db.get_file_count()
        self._status(f"Phase 1 done: {count:,} files in {p1_time:.1f}s — search ready!")
        self._progress(count)
        self._phase(1, "done")
        self.db.set_meta("last_index_time", datetime.now().isoformat())

        # ═══ PHASE 2: Content indexing (background, parallel) ═══
        if self._stop_event.is_set():
            return
        self._phase(2, "starting")
        self._status("Phase 2: Indexing file content...")
        p2_count = self._phase2_content()
        p2_time = time.time() - start_time - p1_time

        if self._stop_event.is_set():
            return

        self._status(f"Phase 2 done: {p2_count:,} files enriched in {p2_time:.1f}s")
        self._phase(2, "done")

        # ═══ PHASE 3: Hash computation (background, parallel) ═══
        if self._stop_event.is_set():
            return
        self._phase(3, "starting")
        self._status("Phase 3: Computing file hashes...")
        p3_count = self._phase3_hash()
        p3_time = time.time() - start_time - p1_time - p2_time

        if self._stop_event.is_set():
            return

        self._phase(3, "done")

        # ═══ Finalize ═══
        elapsed = time.time() - start_time
        self.db.set_meta("last_index_time", datetime.now().isoformat())
        self.db.set_meta("index_duration", f"{elapsed:.1f}")
        self.db.set_meta("total_files", str(count))

        if reindex:
            try:
                with self.db.lock:
                    self.db.conn.execute("PRAGMA wal_checkpoint(TRUNCATE)")
                    self.db.conn.execute("PRAGMA journal_mode=DELETE")
                    self.db.conn.execute("VACUUM")
                    self.db.conn.execute("PRAGMA journal_mode=WAL")
                    self.db.conn.execute("PRAGMA wal_autocheckpoint=200")
            except Exception:
                pass
        else:
            try:
                with self.db.lock:
                    self.db.conn.execute("PRAGMA wal_checkpoint(TRUNCATE)")
            except Exception:
                pass

        count = self.db.get_file_count()
        db_size = self.db.get_db_size_mb()
        self._status(f"Ready! {count:,} files ({db_size:.0f} MB) — {elapsed:.1f}s total")
        self._progress(count)

    # ─── Phase 1: Metadata Scan ──────────────────────────

    def _phase1_scan(self, drives):
        """Scan all drives in parallel — metadata only, no content/hash reading."""
        batch = []
        batch_size = 2000

        # Use ThreadPoolExecutor to scan drives in parallel
        with ThreadPoolExecutor(max_workers=min(len(drives), 4)) as executor:
            futures = {}
            for drive in drives:
                if self._stop_event.is_set():
                    break
                future = executor.submit(self._scan_drive_metadata, drive)
                futures[future] = drive

            for future in as_completed(futures):
                if self._stop_event.is_set():
                    break
                drive = futures[future]
                try:
                    results = future.result()
                    if results:
                        batch.extend(results)
                        # Flush in chunks
                        while len(batch) >= batch_size:
                            chunk = batch[:batch_size]
                            batch = batch[batch_size:]
                            if self._is_reindex:
                                self.db.insert_batch(chunk)
                            else:
                                self.db.upsert_batch(chunk)
                            self.total_indexed += len(chunk)
                            self._progress(self.total_indexed)
                            self._status(f"Phase 1: {self.total_indexed:,} files scanned...")
                except Exception:
                    pass

        # Flush remaining
        if batch:
            if self._is_reindex:
                self.db.insert_batch(batch)
            else:
                self.db.upsert_batch(batch)
            self.total_indexed += len(batch)
            self._progress(self.total_indexed)

    def _scan_drive_metadata(self, drive):
        """Scan a single drive — returns list of metadata tuples. Runs in thread."""
        results = []
        self._scan_dir_recursive(drive, results)
        return results

    def _scan_dir_recursive(self, root, results):
        """Recursively scan directory, collecting metadata only."""
        try:
            for entry in os.scandir(root):
                if self._stop_event.is_set():
                    return
                try:
                    name = entry.name
                    if entry.is_dir(follow_symlinks=False):
                        if name in self.SKIP_DIRS or name.startswith('.'):
                            continue
                        try:
                            self._scan_dir_recursive(entry.path, results)
                        except (PermissionError, OSError):
                            pass
                    else:
                        ext = Path(name).suffix.lower()
                        if ext in self.SKIP_EXTENSIONS:
                            continue
                        try:
                            stat = entry.stat()
                            size = stat.st_size
                            modified = stat.st_mtime
                        except (OSError, PermissionError):
                            size = 0
                            modified = 0

                        # Incremental: skip unchanged files
                        if not self._is_reindex and self._incremental_since > 0 and modified <= self._incremental_since:
                            continue

                        folder_name = os.path.basename(root)
                        dir_path = root.rstrip("\\")

                        # Phase 1: metadata only — no content, no hash
                        results.append((name, dir_path, ext, size, modified, folder_name))

                except (PermissionError, OSError):
                    continue
        except (PermissionError, OSError):
            pass

    # ─── Phase 2: Content Indexing ───────────────────────

    def _phase2_content(self):
        """Read content for text/rich files using thread pool."""
        total_enriched = 0

        while not self._stop_event.is_set():
            files = self.db.get_files_needing_content(limit=500)
            if not files:
                break

            updates = []
            with ThreadPoolExecutor(max_workers=4) as executor:
                futures = {}
                for name, dir_path, ext, size in files:
                    if self._stop_event.is_set():
                        break
                    path = os.path.join(dir_path, name)
                    future = executor.submit(self.read_content, path, ext, size)
                    futures[future] = (dir_path, name)

                for future in as_completed(futures):
                    if self._stop_event.is_set():
                        break
                    dir_path, name = futures[future]
                    try:
                        content = future.result(timeout=10)  # 10s timeout per file
                        # ALWAYS mark as attempted — " " for empty/failed so it won't retry
                        updates.append((content if content else " ", dir_path, name))
                    except Exception:
                        updates.append((" ", dir_path, name))

            if updates:
                self.db.update_content_batch(updates)
                total_enriched += len(updates)
                self._status(f"Phase 2: {total_enriched:,} files enriched...")

        return total_enriched

    # ─── Phase 3: Hash Computation ───────────────────────

    def _phase3_hash(self):
        """Compute hashes for all files using thread pool."""
        total_hashed = 0

        while not self._stop_event.is_set():
            files = self.db.get_files_needing_hash(limit=2000)
            if not files:
                break

            updates = []
            with ThreadPoolExecutor(max_workers=6) as executor:
                futures = {}
                for name, dir_path, size in files:
                    if self._stop_event.is_set():
                        break
                    path = os.path.join(dir_path, name)
                    future = executor.submit(self.compute_file_hash, path, size)
                    futures[future] = (dir_path, name)

                for future in as_completed(futures):
                    if self._stop_event.is_set():
                        break
                    dir_path, name = futures[future]
                    try:
                        file_hash = future.result(timeout=5)  # 5s timeout
                        updates.append((file_hash if file_hash else "-", dir_path, name))
                    except Exception:
                        updates.append(("-", dir_path, name))

            if updates:
                self.db.update_hash_batch(updates)
                total_hashed += len(updates)
                self._status(f"Phase 3: {total_hashed:,} files hashed...")

        return total_hashed


# ══════════════════════════════════════════════════════════════
#  FILE WATCHER
# ══════════════════════════════════════════════════════════════

class QuickFindEventHandler(FileSystemEventHandler):
    def __init__(self, db: FileDatabase, status_callback=None):
        super().__init__()
        self.db = db
        self.status_callback = status_callback

    def _should_skip(self, path):
        parts = path.replace("/", "\\").split("\\")
        for p in parts:
            if p in FileIndexer.SKIP_DIRS or p.startswith('.'):
                return True
        ext = Path(path).suffix.lower()
        return ext in FileIndexer.SKIP_EXTENSIONS

    def _status(self, msg):
        if self.status_callback:
            self.status_callback(msg)

    def _index_single(self, path):
        if self._should_skip(path):
            return
        try:
            if os.path.isdir(path):
                return
            name = os.path.basename(path)
            ext = Path(name).suffix.lower()
            try:
                stat = os.stat(path)
                size = stat.st_size
                modified = stat.st_mtime
            except OSError:
                size = 0
                modified = 0
            dir_path = os.path.dirname(path)
            folder_name = os.path.basename(dir_path)
            content = FileIndexer.read_content(path, ext, size)
            file_hash = FileIndexer.compute_file_hash(path, size) if size > 0 else ""
            self.db.upsert_file(name, dir_path, ext, size, modified, folder_name, content, file_hash)
        except Exception:
            pass

    def on_created(self, event):
        if self._should_skip(event.src_path):
            return
        self._index_single(event.src_path)
        self._status(f"Added: {os.path.basename(event.src_path)}")

    def on_deleted(self, event):
        if self._should_skip(event.src_path):
            return
        if event.is_directory:
            self.db.delete_by_dir_prefix(event.src_path)
        else:
            self.db.delete_by_path(event.src_path)
        self._status(f"Removed: {os.path.basename(event.src_path)}")

    def on_modified(self, event):
        if event.is_directory or self._should_skip(event.src_path):
            return
        self._index_single(event.src_path)

    def on_moved(self, event):
        if event.is_directory:
            self.db.delete_by_dir_prefix(event.src_path)
        else:
            self.db.delete_by_path(event.src_path)
        if not self._should_skip(event.dest_path):
            self._index_single(event.dest_path)
            self._status(f"Moved: {os.path.basename(event.dest_path)}")


class FileWatcher:
    def __init__(self, db: FileDatabase, status_callback=None):
        self.db = db
        self.handler = QuickFindEventHandler(db, status_callback)
        self.observer = Observer()
        self._running = False
        self._heartbeat_timer = None

    def start(self):
        if self._running:
            return
        for drive in _detect_drives():
            try:
                self.observer.schedule(self.handler, drive, recursive=True)
            except Exception:
                pass
        self.observer.daemon = True
        self.observer.start()
        self._running = True
        self._update_heartbeat()
        self._heartbeat_timer = threading.Timer(30, self._heartbeat_loop)
        self._heartbeat_timer.daemon = True
        self._heartbeat_timer.start()

    def _update_heartbeat(self):
        try:
            self.db.set_meta("watcher_heartbeat", datetime.now().isoformat())
        except Exception:
            pass

    def _heartbeat_loop(self):
        while self._running:
            self._update_heartbeat()
            time.sleep(30)

    def stop(self):
        if self._running:
            self._running = False
            self.observer.stop()
            self.observer.join(timeout=5)
            try:
                self.db.set_meta("watcher_heartbeat", "")
            except Exception:
                pass

    def is_running(self):
        return self._running

    @staticmethod
    def is_watcher_active(db):
        hb = db.get_meta("watcher_heartbeat")
        if not hb:
            return False
        try:
            dt = datetime.fromisoformat(hb)
            return (datetime.now() - dt).total_seconds() < 120
        except Exception:
            return False
